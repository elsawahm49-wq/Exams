from pptx import Presentation
from docx import Document
from pypdf import PdfReader

def from_pdf(path: str) -> str:
    r = PdfReader(path)
    return "\n".join([p.extract_text() or "" for p in r.pages])

def from_pptx(path: str) -> str:
    prs = Presentation(path)
    texts = []
    for slide in prs.slides:
        for shp in slide.shapes:
            if hasattr(shp, "text"):
                texts.append(shp.text)
    return "\n".join(texts)

def from_docx(path: str) -> str:
    doc = Document(path)
    return "\n".join(p.text for p in doc.paragraphs
